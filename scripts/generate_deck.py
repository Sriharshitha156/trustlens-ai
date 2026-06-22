#!/usr/bin/env python3
"""Generate the TrustLens AI design-rationale slide deck as editable .pptx files.

Run:  python scripts/generate_deck.py
Output:
  docs/TrustLens-AI-Deck.pptx        (12 slides, 16:9, dark theme)
  docs/TrustLens-AI-Deck-Light.pptx  (12 slides, 16:9, light theme)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def C(hexstr):
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


# ---- two palettes ----
DARK = dict(
    BG=C("0B1220"), CARD=C("111C30"), BLUE=C("0076CE"), BLUE2=C("1DA1F2"),
    TXT=C("F8FAFC"), SEC=C("CBD5E1"), MUT=C("94A3B8"), OK=C("22C55E"),
    WARN=C("F59E0B"), LINE=C("27384F"), HEADER_BG=C("16243A"), HEADER_TXT=C("F8FAFC"),
)
LIGHT = dict(
    BG=C("FFFFFF"), CARD=C("F1F5F9"), BLUE=C("0076CE"), BLUE2=C("0A66C2"),
    TXT=C("0F172A"), SEC=C("334155"), MUT=C("64748B"), OK=C("15803D"),
    WARN=C("B45309"), LINE=C("CBD5E1"), HEADER_BG=C("E2E8F0"), HEADER_TXT=C("0F172A"),
)

ROOT = Path(__file__).resolve().parent.parent


def build_deck(P, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    BG, CARD, BLUE, BLUE2 = P["BG"], P["CARD"], P["BLUE"], P["BLUE2"]
    TXT, SEC, MUT = P["TXT"], P["SEC"], P["MUT"]
    OK, WARN, LINE = P["OK"], P["WARN"], P["LINE"]
    HEADER_BG, HEADER_TXT = P["HEADER_BG"], P["HEADER_TXT"]

    def slide():
        s = prs.slides.add_slide(BLANK)
        bg = s.shapes.add_shape(1, 0, 0, SW, SH)
        bg.fill.solid(); bg.fill.fore_color.rgb = BG
        bg.line.fill.background(); bg.shadow.inherit = False
        bar = s.shapes.add_shape(1, 0, 0, SW, Pt(6))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background(); bar.shadow.inherit = False
        return s

    def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        box = s.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        return tf

    def para(tf, text, size, color, bold=False, italic=False, space_after=8,
             align=PP_ALIGN.LEFT, first=False):
        p = tf.paragraphs[0] if first and tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        run = p.add_run(); run.text = text
        f = run.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = "Segoe UI"
        return p

    def kicker(s, text):
        tf = tb(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5))
        para(tf, text.upper(), 14, BLUE2, bold=True, first=True)

    def heading(s, text, size=34, top=1.05):
        tf = tb(s, Inches(0.9), Inches(top), Inches(11.5), Inches(1.3))
        para(tf, text, size, TXT, bold=True, first=True)

    def card(s, l, t, w, h, title, lines):
        box = s.shapes.add_shape(1, l, t, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE; box.line.width = Pt(1); box.shadow.inherit = False
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25); tf.margin_top = Inches(0.18)
        para(tf, title, 18, BLUE2, bold=True, first=True, space_after=6)
        for ln in lines:
            para(tf, ln, 15, SEC, space_after=4)

    def bullets(s, items, top=2.3, size=20, left=0.95, width=11.4):
        tf = tb(s, Inches(left), Inches(top), Inches(width), Inches(4.6))
        for i, it in enumerate(items):
            para(tf, "•  " + it, size, SEC, first=(i == 0), space_after=12)

    def table(s, rows, top=2.3, left=0.95, width=11.4, col_w=None, fontsize=15):
        nrows = len(rows); ncols = len(rows[0])
        h = Inches(0.5 * nrows + 0.2)
        gt = s.shapes.add_table(nrows, ncols, Inches(left), Inches(top), Inches(width), h).table
        if col_w:
            for i, cw in enumerate(col_w):
                gt.columns[i].width = Inches(cw)
        for r in range(nrows):
            for c in range(ncols):
                cell = gt.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG if r == 0 else CARD
                cell.margin_left = Inches(0.12); cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)
                p = cell.text_frame.paragraphs[0]
                cell.text_frame.word_wrap = True
                run = p.add_run(); run.text = str(rows[r][c])
                run.font.size = Pt(fontsize)
                run.font.bold = (r == 0)
                run.font.color.rgb = HEADER_TXT if r == 0 else SEC
                run.font.name = "Segoe UI"
        return gt

    def metrics(s, items, top=2.6):
        n = len(items); gap = Inches(0.4); total_w = Inches(11.4)
        w = Emu(int((total_w - gap * (n - 1)) / n)); x = Inches(0.95)
        for (num, label) in items:
            box = s.shapes.add_shape(1, x, Inches(top), w, Inches(1.7))
            box.fill.solid(); box.fill.fore_color.rgb = CARD
            box.line.color.rgb = LINE; box.shadow.inherit = False
            tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            para(tf, num, 40, BLUE2, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=2)
            para(tf, label, 14, MUT, align=PP_ALIGN.CENTER)
            x = Emu(int(x) + int(w) + int(gap))

    # ---------- SLIDE 1 ----------
    s = slide()
    kicker(s, "Dell AI Future Minds Hackathon 2026 · Team AI Mavericks")
    tf = tb(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(3))
    para(tf, "TrustLens AI", 60, TXT, bold=True, first=True, space_after=10)
    para(tf, "Transparent & Trustworthy AI Agent Interfaces for IT Operations", 24, MUT, space_after=20)
    para(tf, "“We turn opaque AI recommendations into decisions an IT admin can understand, question, and stand behind.”",
         22, TXT, bold=True)

    # ---------- SLIDE 2 ----------
    s = slide()
    kicker(s, "The Problem")
    heading(s, "AI agents act — but their interfaces are opaque")
    card(s, Inches(0.95), Inches(2.2), Inches(5.5), Inches(1.9), "Opacity",
         ["No view into why a recommendation was made — what data, what logic."])
    card(s, Inches(6.85), Inches(2.2), Inches(5.5), Inches(1.9), "Calibration uncertainty",
         ["“87%” tells you nothing about whether to trust it or question it."])
    card(s, Inches(0.95), Inches(4.3), Inches(5.5), Inches(1.9), "Accountability gap",
         ["When AI is wrong, there's no usable audit trail to course-correct."])
    card(s, Inches(6.85), Inches(4.3), Inches(5.5), Inches(1.9), "Result: a trust deficit",
         ["Hesitation, manual overrides, or avoidance — killing the value AI should deliver."])

    # ---------- SLIDE 3 ----------
    s = slide()
    kicker(s, "Who We Designed For")
    heading(s, "Three personas, one platform")
    table(s, [
        ["Persona", "Needs", "What we gave them"],
        ["IT Admin (primary)", "Speed, accuracy, audit traceability", "Governance dashboard, Trust Analytics, Audit Center, override authority"],
        ["Security Analyst", "Evidence + confidence to escalate", "Evidence weighting, similar cases, outcome learning"],
        ["Non-tech stakeholder", "Plain-language “what & why”", "Simplified view, plain audit log"],
    ], top=2.2, col_w=[3.1, 3.6, 4.7], fontsize=14)
    tf = tb(s, Inches(0.95), Inches(5.7), Inches(11.4), Inches(1))
    para(tf, "Key insight: same data, three lenses — the UI is role-gated, not three separate apps.",
         18, TXT, bold=True, first=True)

    # ---------- SLIDE 4 ----------
    s = slide()
    kicker(s, "Design Principles")
    heading(s, "How we build calibrated trust")
    bullets(s, [
        "No ML jargon, ever. Confidence is a band (“High”, “Review Recommended”), never a raw probability.",
        "Show the reasoning, not just the verdict. Every card answers why.",
        "Human-in-the-loop by default. No consequential action without explicit confirmation.",
        "Make doubt visible. Limitations and counter-considerations are first-class.",
        "Calibrated trust > blind trust. History sits next to every confidence score.",
    ], top=2.2)

    # ---------- SLIDE 5 ----------
    s = slide()
    kicker(s, "Core Requirement")
    heading(s, "The five mandatory transparency elements")
    table(s, [
        ["Element", "In TrustLens AI"],
        ["Reasoning steps", "7-step “Why Did The AI Recommend This?” + evidence weighting"],
        ["Confidence level", "Qualitative bands + plain explanation (no raw %)"],
        ["Data source attribution", "Source cards with trust levels"],
        ["Known limitations", "AI Limitation Awareness + Counter Consideration + ledger weaknesses"],
        ["Human-in-the-loop", "Approve / Override / Escalate + Adaptive Approval Gate"],
    ], top=2.1, col_w=[3.3, 8.1], fontsize=14)
    tf = tb(s, Inches(0.95), Inches(6.4), Inches(11.4), Inches(0.6))
    para(tf, "✓ All five are present and mandatory in the prototype.", 17, OK, bold=True, first=True)

    # ---------- SLIDE 6 ----------
    s = slide()
    kicker(s, "Signature Pattern #1")
    heading(s, "Calibrated confidence, not a raw number")
    card(s, Inches(0.95), Inches(2.3), Inches(5.5), Inches(1.7), "Before",
         ["“87%” — meaningless on its own."])
    card(s, Inches(6.85), Inches(2.3), Inches(5.5), Inches(1.7), "After",
         ["“Review Recommended — based on telemetry from 342 similar devices over 14 days.”"])
    tf = tb(s, Inches(0.95), Inches(4.4), Inches(11.4), Inches(2))
    para(tf, "Paired with the Trust Ledger: “Correct 47 / Incorrect 3 · over-flags Linux devices.”",
         20, SEC, first=True, space_after=8)
    para(tf, "Users calibrate from history + scope, not from a number alone.", 20, TXT, bold=True)

    # ---------- SLIDE 7 ----------
    s = slide()
    kicker(s, "Signature Pattern #2")
    heading(s, "Adaptive Approval Gate")
    tf = tb(s, Inches(0.95), Inches(2.0), Inches(11.4), Inches(0.6))
    para(tf, "Friction scales with risk — low-risk isn't annoying, high-risk isn't reckless.",
         20, MUT, first=True)
    table(s, [
        ["Risk", "Gate"],
        ["Low", "One-click"],
        ["Medium", "Evidence review"],
        ["High", "Evidence + impact review"],
        ["Critical", "+ written justification"],
    ], top=2.7, col_w=[3.0, 8.4], fontsize=16)
    tf = tb(s, Inches(0.95), Inches(6.2), Inches(11.4), Inches(0.9))
    para(tf, "Why: uniform friction trains users to click through everything; adaptive friction preserves attention for what matters.",
         16, TXT, bold=True, first=True)

    # ---------- SLIDE 8 ----------
    s = slide()
    kicker(s, "Signature Pattern #3")
    heading(s, "Designing for doubt")
    card(s, Inches(0.95), Inches(2.2), Inches(5.5), Inches(1.9), "Counter Consideration",
         ["An alternative explanation shown next to the recommendation — fights automation bias."])
    card(s, Inches(6.85), Inches(2.2), Inches(5.5), Inches(1.9), "Impact Preview",
         ["Consequences of both approving and dismissing."])
    card(s, Inches(0.95), Inches(4.3), Inches(5.5), Inches(1.9), "AI Incident Cards",
         ["When AI is wrong: a structured “what failed & lesson learned.”"])
    card(s, Inches(6.85), Inches(4.3), Inches(5.5), Inches(1.9), "AI Limitation Awareness",
         ["Flagged telemetry gaps with impact levels."])

    # ---------- SLIDE 9 ----------
    s = slide()
    kicker(s, "Innovation & Stretch Goals")
    heading(s, "Beyond the brief")
    bullets(s, [
        "High-Risk Decision Courtroom — Advocate vs Challenger agents debate the SAME evidence; a disagreement meter forces human review on contested cases.",
        "Multi-Agent Transparency — visualizes the Detection → Analysis → Remediation handoff.",
        "Problem-first framing + Priority-triage browser (risk + impact + aging + confidence).",
        "Accessibility (WCAG 2.1 AA) — Lighthouse 100/100; merged Audit & Activity Log with CSV export.",
    ], top=2.2)

    # ---------- SLIDE 10 ----------
    s = slide()
    kicker(s, "Validation")
    heading(s, "We tested with real users")
    tf = tb(s, Inches(0.95), Inches(2.0), Inches(11.4), Inches(0.6))
    para(tf, "5 moderated think-aloud sessions · non-ML participants", 20, MUT, first=True)
    metrics(s, [("5", "participants (meets minimum)"),
                ("64%", "comprehension (“what & why”)"),
                ("100", "Lighthouse accessibility")], top=2.7)
    tf = tb(s, Inches(0.95), Inches(4.8), Inches(11.4), Inches(2))
    para(tf, "Found & fixing: first-time orientation · confidence-vs-reliability clarity · control discoverability.",
         18, SEC, first=True, space_after=6)
    para(tf, "“In control” ~50% → our top improvement.", 18, WARN, bold=True, space_after=10)
    para(tf, "Full method & results in our Usability Test Report.", 15, MUT, italic=True)

    # ---------- SLIDE 11 ----------
    s = slide()
    kicker(s, "Architecture — Honest to the Brief")
    heading(s, "Simulated AI, real UX focus")
    bullets(s, [
        "No model trained, no prod backend — outputs simulated from Faker-generated IT logs (16 linked CSVs).",
        "React 19 + Vite + Tailwind, fully client-side; loads CSVs via PapaParse.",
        "Runs offline — no API key required.",
        "60%+ of effort on design + testing, per the brief's guidance.",
    ], top=2.2)

    # ---------- SLIDE 12 ----------
    s = slide()
    kicker(s, "Impact & Next Steps")
    heading(s, "Calibrated trust, by design")
    tf = tb(s, Inches(0.95), Inches(2.2), Inches(11.4), Inches(3))
    para(tf, "Admins lean on AI where it's strong, override where they know better — and always have a record.",
         24, TXT, bold=True, first=True, space_after=18)
    para(tf, "Next: live Hugging Face inference + SHAP rendering · multi-tenant audit retention · forced first-run onboarding.",
         18, SEC, space_after=24)
    para(tf, "Thank you — live demo now.", 22, BLUE2, bold=True)

    prs.save(str(out_path))
    print(f"Saved -> {out_path}")


if __name__ == "__main__":
    build_deck(DARK, ROOT / "docs" / "TrustLens-AI-Deck.pptx")
    build_deck(LIGHT, ROOT / "docs" / "TrustLens-AI-Deck-Light.pptx")
